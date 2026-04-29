#!/usr/bin/env python3
"""
Generate a PowerPoint deck from emails matching a Gmail search query.

CLI: python ppt.py "<gmail-query>"
Example: python ppt.py "from:investor newer_than:7d"

Pipeline:
  1. Fetch up to 25 emails matching the query.
  2. Summarize each via local Ollama (JSON-mode).
  3. Generate cross-email "Top 3" summary.
  4. Build .pptx and save to ~/email-ppts/YYYY-MM-DD-HHMMSS.pptx.
  5. Send Telegram completion ping.
"""

import os
import sys
import json
import logging
import traceback
from pathlib import Path
from datetime import datetime, timezone

from firestore_activity import get_db, report_run
from firestore_alerts import send_alert
from firestore_users import load_user_config
from mime_extract import extract_body, decode_header_value
from lang_hint import detect_dominant_script, language_directive

from dotenv import load_dotenv
from openai import OpenAI
from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request
from googleapiclient.discovery import build

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------- Config ----------
BASE_DIR = Path(__file__).parent.resolve()
OUTPUT_DIR = Path.home() / "email-ppts"

load_dotenv(BASE_DIR / ".env")

TELEGRAM_BOT_TOKEN = os.environ["TELEGRAM_BOT_TOKEN"]
AUTHORIZED_CHAT_ID = int(os.environ["AUTHORIZED_CHAT_ID"])
USER_UID = os.environ.get("EMAIL2PPT_USER_UID", "").strip()
OLLAMA_BASE_URL = os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434/v1")
OLLAMA_MODEL = os.environ.get("OLLAMA_MODEL", "llama3.1:8b")

GMAIL_TOKEN_PATH = BASE_DIR / "gmail_token.json"
GMAIL_SCOPES = ["https://www.googleapis.com/auth/gmail.modify"]

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

LOG_PATH = BASE_DIR / "ppt.log"
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    handlers=[logging.FileHandler(LOG_PATH), logging.StreamHandler()],
)
log = logging.getLogger("ppt")
from log_redaction import install_redaction_filter  # noqa: E402
install_redaction_filter(logging.getLogger())


# ---------- Gmail helpers ----------
def get_gmail_service():
    creds = Credentials.from_authorized_user_file(str(GMAIL_TOKEN_PATH), GMAIL_SCOPES)
    if creds.expired and creds.refresh_token:
        creds.refresh(Request())
        # Refresh tokens grant gmail.modify indefinitely; ensure 0600 every write.
        GMAIL_TOKEN_PATH.write_text(creds.to_json())
        os.chmod(GMAIL_TOKEN_PATH, 0o600)
    return build("gmail", "v1", credentials=creds)


def fetch_emails(query: str, max_results: int = 25) -> list:
    service = get_gmail_service()
    log.info(f"Gmail query: {query}")
    result = (
        service.users()
        .messages()
        .list(userId="me", q=query, maxResults=max_results)
        .execute()
    )
    messages = result.get("messages", [])
    emails = []
    for m in messages:
        msg = (
            service.users()
            .messages()
            .get(userId="me", id=m["id"], format="full")
            .execute()
        )
        headers = {h["name"]: h["value"] for h in msg["payload"]["headers"]}
        body = extract_body(msg["payload"])[:4000]
        emails.append(
            {
                "from": decode_header_value(headers.get("From", "")),
                "subject": decode_header_value(headers.get("Subject", "")),
                "date": headers.get("Date", ""),
                "snippet": msg.get("snippet", ""),
                "body": body,
            }
        )
    return emails


# ---------- LLM summarization ----------
DEFAULT_PERSONA_LINE = "You are an executive assistant summarizing emails for the recipient."

EMAIL_USER_TEMPLATE = """Summarize the email below into structured JSON ONLY (no prose, no markdown fences).

Required JSON shape:
{{
  "context": [1-2 short strings describing who the sender is and why they are writing],
  "key_points": [up to {kp_max} short strings, each one specific fact, claim, request, number, or quote from the email],
  "asks": [up to {asks_max} short strings describing what the sender wants from the recipient; if no action is requested, return a single string stating that],
  "suggested_response": "one short string describing what to do next, or that no reply is needed",
  "urgency": "low" | "med" | "high"
}}

Style:
- Specific over generic. Pull real numbers, names, dates, dollar figures.
- Skip greetings, signatures, quoted-thread history.
- Output ONLY the JSON object.

EMAIL:
From: {from_}
Subject: {subject}
Date: {date}

{body}
"""


def _build_system_message(cfg: dict) -> str:
    extra = (cfg.get("summaryPersona") or "").strip()
    if extra:
        return f"{DEFAULT_PERSONA_LINE}\n\nAdditional instructions: {extra}"
    return DEFAULT_PERSONA_LINE


def _strip_fences(text: str) -> str:
    text = text.strip()
    if text.startswith("```"):
        parts = text.split("```")
        if len(parts) >= 2:
            text = parts[1]
            if text.startswith("json"):
                text = text[4:]
            text = text.strip()
    return text


def summarize_email(client: OpenAI, email: dict, cfg: dict, system_msg: str) -> dict:
    user_msg = EMAIL_USER_TEMPLATE.format(
        kp_max=cfg.get("summaryKeyPointsMax", 6),
        asks_max=cfg.get("summaryAsksMax", 4),
        from_=email["from"],
        subject=email["subject"],
        date=email["date"],
        body=email["body"],
    )
    directive = language_directive(detect_dominant_script(email["body"]))
    if directive:
        user_msg = f"{user_msg}\n\n{directive}"
    resp = client.chat.completions.create(
        model=OLLAMA_MODEL,
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_msg},
        ],
        temperature=0.3,
    )
    text = _strip_fences(resp.choices[0].message.content or "")
    try:
        data = json.loads(text)
    except json.JSONDecodeError:
        log.warning(f"JSON parse failed for {email['subject']!r}; using fallback")
        data = {
            "context": [],
            "key_points": [text[:400]] if text else [],
            "asks": [],
            "suggested_response": "",
            "urgency": "low",
        }
    return data


def top3_summary(client: OpenAI, summarized: list, system_msg: str) -> list:
    blocks = []
    for s in summarized:
        kp = "; ".join(s["data"].get("key_points", [])[:3])
        blocks.append(f"- {s['email']['from']}: {s['email']['subject']} | key: {kp}")
    user_msg = (
        "Given the email summaries below, return EXACTLY 3 short strings capturing "
        "the most important takeaways across all emails. Output ONLY a JSON array "
        "of 3 strings. No prose, no fences.\n\n" + "\n".join(blocks)
    )
    resp = client.chat.completions.create(
        model=OLLAMA_MODEL,
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_msg},
        ],
        temperature=0.3,
    )
    text = _strip_fences(resp.choices[0].message.content or "")
    try:
        items = json.loads(text)
        if isinstance(items, list) and items:
            return [str(x) for x in items[:3]]
    except json.JSONDecodeError:
        pass
    return [s["email"]["subject"] or s["email"]["from"] for s in summarized[:3]]


# ---------- PPTX builder ----------
URGENCY_COLOR = {
    "low": RGBColor(0x34, 0xC7, 0x59),
    "med": RGBColor(0xFF, 0x9F, 0x0A),
    "high": RGBColor(0xFF, 0x3B, 0x30),
}


def _add_text(
    slide, text, left, top, width, height, size,
    bold=False, color=None, align=None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return tb


def _short_sender(from_field: str) -> str:
    if "<" in from_field:
        name = from_field.split("<")[0].strip().strip('"')
        return name or from_field
    return from_field


def build_deck(query: str, top3: list, summarized: list, run_at: datetime) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Title
    s1 = prs.slides.add_slide(blank)
    _add_text(s1, "Email Digest", Inches(0.6), Inches(2.4), Inches(12), Inches(1.4), 54, bold=True)
    _add_text(s1, run_at.strftime("%A, %B %d, %Y"), Inches(0.6), Inches(3.8), Inches(12), Inches(0.6), 24)
    _add_text(
        s1, f"Query: {query}",
        Inches(0.6), Inches(4.5), Inches(12), Inches(0.6),
        16, color=RGBColor(0x86, 0x86, 0x8A),
    )

    # Top 3
    s2 = prs.slides.add_slide(blank)
    _add_text(s2, "Top 3 Things to Know", Inches(0.6), Inches(0.5), Inches(12), Inches(1), 36, bold=True)
    for i, item in enumerate(top3):
        _add_text(s2, f"• {item}", Inches(0.8), Inches(2.0 + 0.9 * i), Inches(12), Inches(1), 22)

    # Per email
    for s in summarized:
        e = s["email"]
        d = s["data"]
        slide = prs.slides.add_slide(blank)

        _add_text(slide, _short_sender(e["from"]), Inches(0.6), Inches(0.4), Inches(11), Inches(0.7), 28, bold=True)
        _add_text(slide, e["subject"], Inches(0.6), Inches(1.0), Inches(11), Inches(0.6), 16, color=RGBColor(0x48, 0x48, 0x4A))

        urgency = (d.get("urgency") or "low").lower()
        if urgency not in URGENCY_COLOR:
            urgency = "low"
        _add_text(
            slide, urgency.upper(),
            Inches(11.5), Inches(0.4), Inches(1.3), Inches(0.5),
            14, bold=True, color=URGENCY_COLOR[urgency], align=PP_ALIGN.RIGHT,
        )

        y = 1.8
        for label, key, max_items in [
            ("Context", "context", 2),
            ("Key Points", "key_points", 6),
            ("Asks / Actions", "asks", 4),
        ]:
            items = d.get(key) or []
            if not items:
                continue
            _add_text(slide, label, Inches(0.6), Inches(y), Inches(11), Inches(0.4), 16, bold=True)
            y += 0.45
            for b in items[:max_items]:
                _add_text(slide, f"• {b}", Inches(0.8), Inches(y), Inches(11), Inches(0.4), 14)
                y += 0.35
            y += 0.2

        if d.get("suggested_response"):
            _add_text(
                slide, f"→ {d['suggested_response']}",
                Inches(0.6), Inches(6.6), Inches(12), Inches(0.6),
                14, bold=True, color=RGBColor(0x00, 0x71, 0xE3),
            )

    return prs


# ---------- Telegram ----------
# Routes via firestore_alerts: customer's own bot if linked, else env shared bot.
def send_telegram(text: str):
    send_alert(USER_UID, text)


# ---------- Main ----------
def main():
    started = datetime.now(timezone.utc)
    outputs: list[str] = []
    email_count = 0
    try:
        if len(sys.argv) < 2 or not sys.argv[1].strip():
            print('usage: python ppt.py "<gmail-query>"', file=sys.stderr)
            sys.exit(2)

        query = sys.argv[1].strip()
        log.info("=" * 60)
        log.info(f"Starting PPT generation: {query}")
        run_at = datetime.now()

        emails = fetch_emails(query)
        email_count = len(emails)
        log.info(f"Matched emails: {email_count}")

        if not emails:
            send_telegram(f"📭 PPT skipped — no emails matched: {query}")
            report_run("ppt", "ok", started_at=started, email_count=0)
            return

        client = OpenAI(base_url=OLLAMA_BASE_URL, api_key="ollama-local")

        cfg: dict = {}
        if USER_UID:
            try:
                cfg = load_user_config(get_db(), USER_UID)
            except Exception:  # noqa: BLE001 - CLI tool degrades gracefully
                # without a config; persona/limits fall back to defaults rather
                # than aborting the deck build.
                log.warning("uid=%s failed to load user config; using defaults", USER_UID)
        system_msg = _build_system_message(cfg)

        summarized = []
        for e in emails:
            log.info(f"Summarizing: {e['subject'][:60]}")
            data = summarize_email(client, e, cfg, system_msg)
            summarized.append({"email": e, "data": data})

        log.info("Generating Top 3")
        top3 = top3_summary(client, summarized, system_msg)

        prs = build_deck(query, top3, summarized, run_at)
        filename = run_at.strftime("%Y-%m-%d-%H%M%S") + ".pptx"
        out_path = OUTPUT_DIR / filename
        prs.save(str(out_path))
        log.info(f"Saved: {out_path}")
        outputs.append(filename)

        send_telegram(f"📊 PPT saved: {filename} ({email_count} emails)\n{out_path}")
        report_run(
            "ppt",
            "ok",
            started_at=started,
            email_count=email_count,
            outputs=outputs,
        )
    except SystemExit:
        raise
    except Exception:  # noqa: BLE001 - report-then-rethrow guard:
        # we want the activity record written before the process exits.
        report_run(
            "ppt",
            "error",
            started_at=started,
            email_count=email_count,
            outputs=outputs,
            error=traceback.format_exc(),
        )
        raise


if __name__ == "__main__":
    try:
        main()
    except Exception:  # noqa: BLE001 - top-level subprocess guard.
        log.exception("PPT generation failed")
        try:
            # Don't echo the exception text — see digest.py for rationale.
            send_telegram(f"⚠️ PPT failed. See {LOG_PATH}")
        except (OSError, ValueError):
            log.exception("Failed to deliver PPT failure alert")
        raise
